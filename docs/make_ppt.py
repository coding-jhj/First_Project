from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

# ── 컬러 팔레트 (밝고 모던) ──────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INDIGO     = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DK  = RGBColor(0x31, 0x2E, 0x81)
INDIGO_50  = RGBColor(0xEE, 0xF2, 0xFF)
INDIGO_100 = RGBColor(0xE0, 0xE7, 0xFF)
INDIGO_200 = RGBColor(0xC7, 0xD2, 0xFE)
INDIGO_900 = RGBColor(0x1E, 0x1B, 0x4B)
PURPLE     = RGBColor(0x79, 0x3C, 0xEA)
PURPLE_50  = RGBColor(0xF5, 0xF3, 0xFF)
PURPLE_200 = RGBColor(0xDD, 0xD6, 0xFE)
TEXT       = RGBColor(0x1F, 0x29, 0x37)
TEXT_MID   = RGBColor(0x6B, 0x72, 0x80)
TEXT_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
GREEN      = RGBColor(0x05, 0x96, 0x69)
GREEN_BG   = RGBColor(0xD1, 0xFA, 0xE5)
GREEN_200  = RGBColor(0xA7, 0xF3, 0xD0)
RED        = RGBColor(0xDC, 0x26, 0x26)
RED_BG     = RGBColor(0xFE, 0xE2, 0xE2)
RED_200    = RGBColor(0xFC, 0xA5, 0xA5)
AMBER      = RGBColor(0xD9, 0x77, 0x06)
AMBER_BG   = RGBColor(0xFE, 0xF3, 0xC7)
AMBER_200  = RGBColor(0xFD, 0xE6, 0x8A)
GRAY_50    = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_100   = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_200   = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── 헬퍼 함수 ────────────────────────────────────────────────

def set_bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background()

def rect(slide, x, y, w, h, fill=WHITE, lc=None, lw=Pt(0.5)):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.width = lw
    if lc: r.line.color.rgb = lc
    else:  r.line.fill.background()
    return r

def txt(slide, text, x, y, w, h, size=15, bold=False,
        color=TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    r.font.italic = italic
    return tb

def hdr(slide, n, title, sub=""):
    rect(slide, 0, 0, W, Inches(1.85), fill=INDIGO, lc=None)
    rect(slide, 0, Inches(1.85), W, Inches(5.65), fill=WHITE, lc=None)
    if n:
        txt(slide, f"{n:02d}", Inches(12.3), Inches(0.15),
            Inches(0.85), Inches(0.38), size=20, bold=True,
            color=INDIGO_200, align=PP_ALIGN.RIGHT)
    txt(slide, title, Inches(0.65), Inches(0.2),
        Inches(11.5), Inches(0.95), size=36, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.65), Inches(1.3),
            Inches(11.5), Inches(0.42), size=14, color=INDIGO_200)

def card(slide, x, y, w, h, bg=INDIGO_50, lc=INDIGO_100):
    rect(slide, x, y, w, h, fill=bg, lc=lc, lw=Pt(1))

def trow(slide, cells, xs, y, rh=Inches(0.42),
         bg=WHITE, tc=TEXT, bold=False, sz=13):
    for i, (cell, x) in enumerate(zip(cells, xs)):
        w = (xs[i+1]-x) if i+1 < len(xs) else (Inches(12.73)-x)
        rect(slide, x, y, w, rh, fill=bg, lc=GRAY_200, lw=Pt(0.4))
        txt(slide, cell, x+Inches(0.12), y+Inches(0.07),
            w-Inches(0.18), rh-Inches(0.1), size=sz, bold=bold, color=tc)

def thdr(slide, cells, xs, y, rh=Inches(0.44)):
    trow(slide, cells, xs, y, rh=rh, bg=INDIGO, tc=WHITE, bold=True, sz=13)

def callout(slide, text, y=Inches(6.65)):
    rect(slide, Inches(0.6), y, Inches(12.13), Inches(0.65),
         fill=INDIGO_50, lc=INDIGO_100, lw=Pt(1))
    txt(slide, text, Inches(0.85), y+Inches(0.1), Inches(11.6), Inches(0.48),
        size=14, bold=True, color=INDIGO_900, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)

# 왼쪽 인디고 패널
rect(sl, 0, 0, Inches(7.2), H, fill=INDIGO, lc=None)
rect(sl, Inches(7.2), 0, Inches(0.06), H, fill=INDIGO_200, lc=None)

txt(sl, "VoiceGuide",
    Inches(0.75), Inches(1.4), Inches(6.2), Inches(1.5),
    size=54, bold=True, color=WHITE)
txt(sl, "실시간 객체 탐지 기반",
    Inches(0.75), Inches(3.0), Inches(6.2), Inches(0.5),
    size=20, color=INDIGO_200)
txt(sl, "시각장애인 보행 보조 시스템",
    Inches(0.75), Inches(3.52), Inches(6.2), Inches(0.5),
    size=20, color=INDIGO_200)
rect(sl, Inches(0.75), Inches(4.2), Inches(2.2), Pt(2), fill=INDIGO_200, lc=None)
txt(sl, '"스마트폰 카메라 하나로\n장애물을 감지하고, 음성으로 안내합니다."',
    Inches(0.75), Inches(4.35), Inches(6.2), Inches(1.0),
    size=14, color=INDIGO_100, italic=True)
txt(sl, "2026.05",
    Inches(0.75), Inches(6.85), Inches(2.0), Inches(0.38),
    size=13, color=INDIGO_200)

# 오른쪽 팀 정보
txt(sl, "Team", Inches(7.6), Inches(1.2), Inches(5.3), Inches(0.45),
    size=12, color=TEXT_LIGHT, bold=True)
txt(sl, "[팀명]", Inches(7.6), Inches(1.62), Inches(5.3), Inches(0.7),
    size=30, bold=True, color=INDIGO_900)

members = [
    ("[이름 A]", "YOLO 모델 학습"),
    ("[이름 B]", "모바일 추론"),
    ("[이름 C]", "서버 / API"),
    ("[이름 D]", "발표 및 테스트"),
]
for i, (name, role) in enumerate(members):
    y = Inches(2.55) + i * Inches(0.68)
    rect(sl, Inches(7.6), y, Inches(5.1), Inches(0.58),
         fill=INDIGO_50, lc=INDIGO_100, lw=Pt(0.8))
    txt(sl, name, Inches(7.8), y+Inches(0.1), Inches(1.8), Inches(0.35),
        size=14, bold=True, color=INDIGO_900)
    txt(sl, role, Inches(9.75), y+Inches(0.13), Inches(2.7), Inches(0.3),
        size=12, color=TEXT_MID)

badges = [
    ("YOLO11n",        INDIGO),
    ("YOLO26n fallback", GREEN),
    ("ByteTrack",      PURPLE),
    ("Voice AI",       AMBER),
]
for i, (label, color) in enumerate(badges):
    x = Inches(7.6) + i * Inches(1.35)
    rect(sl, x, Inches(5.55), Inches(1.25), Inches(0.48),
         fill=WHITE, lc=color, lw=Pt(1.5))
    txt(sl, label, x+Inches(0.04), Inches(5.6), Inches(1.18), Inches(0.38),
        size=10, bold=True, color=color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — 문제 정의
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 2, "문제 정의", "왜 이 프로젝트가 필요한가?")

# 왼쪽 카드 — 현황
card(sl, Inches(0.6), Inches(2.05), Inches(5.85), Inches(5.1))
rect(sl, Inches(0.6), Inches(2.05), Inches(5.85), Inches(0.5), fill=RED_BG, lc=None)
rect(sl, Inches(0.6), Inches(2.05), Inches(0.07), Inches(0.5), fill=RED, lc=None)
txt(sl, "⚠  현황 & 문제", Inches(0.85), Inches(2.12),
    Inches(5.35), Inches(0.35), size=15, bold=True, color=RED)

stats = [
    ("25만 명",    "국내 시각장애인 수 (2024 보건복지부)"),
    ("보행 중 위험", "돌출 장애물 · 차량 접근 · 낮은 장애물"),
    ("흰 지팡이 한계", "좁은 감지 범위, 원거리 인식 어려움"),
]
for i, (key, val) in enumerate(stats):
    y = Inches(2.7) + i * Inches(1.45)
    rect(sl, Inches(0.82), y, Inches(5.4), Pt(2.5), fill=RED, lc=None)
    txt(sl, key, Inches(0.82), y+Inches(0.12), Inches(5.4), Inches(0.52),
        size=19, bold=True, color=RED)
    txt(sl, val, Inches(0.82), y+Inches(0.65), Inches(5.4), Inches(0.45),
        size=13, color=TEXT_MID)

# 오른쪽 카드 — 해결
card(sl, Inches(6.82), Inches(2.05), Inches(6.1), Inches(5.1))
rect(sl, Inches(6.82), Inches(2.05), Inches(6.1), Inches(0.5), fill=GREEN_BG, lc=None)
rect(sl, Inches(6.82), Inches(2.05), Inches(0.07), Inches(0.5), fill=GREEN, lc=None)
txt(sl, "✓  VoiceGuide 해결", Inches(7.07), Inches(2.12),
    Inches(5.6), Inches(0.35), size=15, bold=True, color=GREEN)

solutions = [
    ("01", "실시간 장애물 감지",
     "스마트폰 카메라로 전방 객체 탐지\n즉각적인 위험 감지"),
    ("02", "모바일 온디바이스 추론",
     "서버 지연 없이 디바이스에서 직접 처리\n오프라인 환경에서도 동작"),
    ("03", "음성 피드백 안내",
     "위험도 기반 직관적 음성 경고\nRisk Score를 자연어로 변환"),
]
for i, (num, title, desc) in enumerate(solutions):
    y = Inches(2.72) + i * Inches(1.45)
    rect(sl, Inches(7.07), y, Inches(0.46), Inches(0.46), fill=INDIGO, lc=None)
    txt(sl, num, Inches(7.07), y+Inches(0.07), Inches(0.46), Inches(0.32),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(7.65), y+Inches(0.02), Inches(4.95), Inches(0.42),
        size=16, bold=True, color=INDIGO_900)
    txt(sl, desc, Inches(7.65), y+Inches(0.48), Inches(4.95), Inches(0.7),
        size=13, color=TEXT_MID)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — 기존 방식의 한계
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 3, "기존 방식의 한계", "왜 기존 솔루션으로는 부족한가")

xs = [Inches(0.6), Inches(3.4), Inches(7.6), Inches(11.1)]
thdr(sl, ["방식", "제공 기능", "한계"], xs, Inches(2.1))
rows = [
    ("흰 지팡이",      "직접 장애물 탐색",    "좁은 범위, 원거리 인식 불가"),
    ("GPS 앱",         "위치 기반 경로 안내",  "실시간 장애물 감지 불가"),
    ("웨어러블 기기",  "일부 센서 활용",       "고비용, 보급률 낮음"),
    ("기존 카메라 앱", "사진 촬영 안내",       "실시간 피드백 없음"),
]
for i, row in enumerate(rows):
    trow(sl, row, xs, Inches(2.54)+i*Inches(0.56),
         bg=WHITE if i%2==0 else GRAY_50)
callout(sl, "핵심:  실시간  +  경량  +  모바일  —  세 조건을 동시에 만족하는 솔루션이 없다")


# ════════════════════════════════════════════════════════════
# SLIDE 4 — 고객여정지도
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 4, "고객여정지도", "Customer Journey Map  —  사용자가 실제로 겪는 Pain Point")

xs4 = [Inches(0.6), Inches(3.5), Inches(7.9), Inches(11.6)]
hdrs4 = [("상황", INDIGO, WHITE), ("Pain Point", RED_BG, RED), ("VoiceGuide 해결", GREEN_BG, GREEN)]
for i, (label, bg_c, tc) in enumerate(hdrs4):
    w = xs4[i+1]-xs4[i]
    rect(sl, xs4[i], Inches(2.1), w, Inches(0.46), fill=bg_c, lc=GRAY_200, lw=Pt(0.5))
    txt(sl, label, xs4[i]+Inches(0.14), Inches(2.17), w-Inches(0.2), Inches(0.33),
        size=13, bold=True, color=tc)

cjm = [
    ("횡단보도 이동",  "차량 접근 인지 어려움",          "객체 탐지 + 음성 경고"),
    ("보도 이동",      "낮은 장애물 미감지",              "지도 패턴 안내"),
    ("혼잡한 장소",    "빠르게 움직이는 객체 추적 불가",  "ByteTrack ID 추적"),
    ("처음 가는 길",   "주변 구조물 인식 불가",           "실시간 음성 안내"),
]
for i, row in enumerate(cjm):
    y_ = Inches(2.56) + i * Inches(0.9)
    for j, (cell, x) in enumerate(zip(row, xs4)):
        w = xs4[j+1]-x if j+1 < len(xs4) else Inches(12.73)-x
        bg_c = WHITE if i%2==0 else GRAY_50
        rect(sl, x, y_, w, Inches(0.82), fill=bg_c, lc=GRAY_200, lw=Pt(0.4))
        tc = TEXT if j==0 else (RED if j==1 else GREEN)
        txt(sl, cell, x+Inches(0.14), y_+Inches(0.16),
            w-Inches(0.22), Inches(0.52), size=13, bold=(j==2), color=tc)

txt(sl, "기술이 아닌 '사용자 상황'에서 출발한 설계",
    Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3),
    size=12, color=TEXT_LIGHT, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — 해결 방향
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 5, "해결 방향", "우리가 선택한 3가지 접근법")

approaches = [
    (INDIGO,  INDIGO_50,  INDIGO_100, "01",
     "모바일 온디바이스 추론",
     ["서버 통신 없이 스마트폰에서 직접 추론",
      "네트워크 지연 없음 / 오프라인 동작"]),
    (PURPLE,  PURPLE_50,  PURPLE_200, "02",
     "적응형 모델 전환 구조",
     ["YOLO11n 기본 탐지 (정확도 중심)",
      "FPS 저하 감지 → YOLO26n 자동 전환"]),
    (GREEN,   GREEN_BG,   GREEN_200,  "03",
     "사용자 중심 피드백",
     ["위험도 Smoothing → 음성 안내",
      "Risk Score를 자연어로 변환"]),
]
for i, (ac, bg_c, bc, num, title, items) in enumerate(approaches):
    x = Inches(0.55) + i * Inches(4.2)
    card(sl, x, Inches(2.05), Inches(3.95), Inches(4.95), bg=bg_c, lc=bc)
    rect(sl, x, Inches(2.05), Inches(3.95), Inches(0.09), fill=ac, lc=None)
    rect(sl, x+Inches(0.2), Inches(2.28), Inches(0.52), Inches(0.52), fill=ac, lc=None)
    txt(sl, num, x+Inches(0.2), Inches(2.31), Inches(0.52), Inches(0.4),
        size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, x+Inches(0.2), Inches(2.96), Inches(3.55), Inches(0.52),
        size=17, bold=True, color=ac if ac != INDIGO else INDIGO_900)
    rect(sl, x+Inches(0.2), Inches(3.54), Inches(3.55), Pt(1.5), fill=ac, lc=None)
    for j, item in enumerate(items):
        txt(sl, f"·  {item}",
            x+Inches(0.2), Inches(3.72)+j*Inches(0.72), Inches(3.6), Inches(0.6),
            size=13, color=TEXT_MID)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — 시스템 아키텍처
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 6, "시스템 아키텍처", "전체 시스템 구성도")

def node(label, x, y, w=Inches(3.0), h=Inches(0.54),
         fill=INDIGO_50, lc=INDIGO_200, tc=INDIGO_900, bold=False):
    rect(sl, x, y, w, h, fill=fill, lc=lc, lw=Pt(1.2))
    txt(sl, label, x+Inches(0.12), y+Inches(0.1), w-Inches(0.18), h-Inches(0.14),
        size=13, bold=bold, color=tc, align=PP_ALIGN.CENTER)

def arr(x, y):
    txt(sl, "▼", x, y, Inches(0.4), Inches(0.3),
        size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

cx = Inches(1.05)
node("📷  스마트폰 카메라",        cx, Inches(2.05), fill=GRAY_100, lc=GRAY_200, tc=TEXT_MID)
arr(Inches(2.4), Inches(2.62))
node("YOLO11n  —  기본 탐지",      cx, Inches(2.97), fill=INDIGO_50, lc=INDIGO_200, bold=True)
arr(Inches(2.4), Inches(3.54))
node("⚡  FPS 저하 감지",           cx, Inches(3.87), fill=AMBER_BG, lc=AMBER_200, tc=AMBER, bold=True)
txt(sl, "↓ 저하 시", Inches(2.25), Inches(4.45), Inches(1.3), Inches(0.3),
    size=11, color=TEXT_LIGHT, italic=True)
node("YOLO26n  —  fallback 전환",  cx, Inches(4.79), fill=GREEN_BG, lc=GREEN_200, tc=GREEN, bold=True)
arr(Inches(2.4), Inches(5.36))
node("ByteTrack  —  객체 추적",    cx, Inches(5.69), fill=INDIGO_50, lc=INDIGO_200, bold=True)

# 하단 출력 분기
rect(sl, Inches(1.05), Inches(6.28), Inches(6.3), Pt(1.5), fill=INDIGO_200, lc=None)
node("위험도 Smoothing",  Inches(1.05), Inches(6.38), w=Inches(2.8))
node("JSON → 서버",       Inches(4.25), Inches(6.38), w=Inches(2.8),
     fill=GRAY_100, lc=GRAY_200, tc=TEXT_MID)

# 오른쪽 설명 패널
card(sl, Inches(7.5), Inches(2.05), Inches(5.45), Inches(5.1))
rect(sl, Inches(7.5), Inches(2.05), Inches(5.45), Inches(0.52), fill=INDIGO, lc=None)
txt(sl, "설계 핵심 포인트", Inches(7.75), Inches(2.13),
    Inches(4.95), Inches(0.35), size=14, bold=True, color=WHITE)

points = [
    ("온디바이스",   "모든 추론이 스마트폰 내에서 실행\n서버 지연 없는 실시간 처리"),
    ("적응형 전환",  "YOLO11n 정상 → YOLO26n fallback\nFPS 저하 시 자동 전환"),
    ("서버 역할",    "로그 저장 및 세션 관리만 담당\nSSOT: JSON 단일 데이터 소스"),
]
for i, (k, v) in enumerate(points):
    y_ = Inches(2.73) + i * Inches(1.42)
    rect(sl, Inches(7.75), y_, Inches(0.07), Inches(0.48), fill=INDIGO_200, lc=None)
    txt(sl, k, Inches(7.98), y_, Inches(4.65), Inches(0.42),
        size=15, bold=True, color=INDIGO_900)
    txt(sl, v, Inches(7.98), y_+Inches(0.44), Inches(4.65), Inches(0.75),
        size=12, color=TEXT_MID)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — 파이프라인 노드
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 7, "파이프라인 노드 설명", "각 단계의 입력 / 처리 / 출력")

xs7 = [Inches(0.6), Inches(3.3), Inches(6.2), Inches(9.6)]
thdr(sl, ["노드", "입력", "처리", "출력"], xs7, Inches(2.1))

pipeline = [
    ("Camera Input",    "—",              "프레임 캡처",             "RGB 이미지"),
    ("YOLO11n",         "이미지",          "객체 탐지 (기본)",         "BBox, Class, Conf"),
    ("FPS Monitor",     "추론 결과",       "FPS 측정 + 임계값 판단",   "Switch 신호"),
    ("YOLO26n",         "이미지",          "객체 탐지 (fallback)",     "BBox, Class, Conf"),
    ("ByteTrack",       "BBox",           "ID 할당 및 추적",           "Track ID, 벡터"),
    ("Risk Analyzer",   "Track + Depth",  "위험도 Smoothing",         "Risk Score (0~1)"),
    ("Feedback Module", "Risk Score",     "임계값 판단 → 음성 생성",  "음성 문자열"),
    ("JSON Sender",     "전체 결과",       "직렬화",                   "JSON → 서버"),
]
highlight = {"YOLO26n": GREEN_BG, "FPS Monitor": AMBER_BG}
for i, row in enumerate(pipeline):
    bg_c = highlight.get(row[0], WHITE if i%2==0 else GRAY_50)
    trow(sl, row, xs7, Inches(2.54)+i*Inches(0.52), bg=bg_c)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — 기술 선택 근거
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 8, "기술 선택 근거", "왜 이 조합인가  —  수치로 설명")

# 왼쪽
card(sl, Inches(0.6), Inches(2.05), Inches(6.0), Inches(5.2))
txt(sl, "YOLO 모델 구조 선택", Inches(0.85), Inches(2.18),
    Inches(5.5), Inches(0.42), size=16, bold=True, color=INDIGO_900)
yx = [Inches(0.65), Inches(2.95), Inches(4.75)]
thdr(sl, ["항목", "YOLO11n", "YOLO26n"], yx, Inches(2.72))
for i, row in enumerate([
    ("역할",       "기본 탐지 모델",  "FPS 저하 fallback"),
    ("우선순위",   "정확도·안정성",   "FPS 안정성"),
    ("경량성",     "경량",            "더 경량"),
    ("특징",       "정밀도 높음",     "빠른 추론 속도"),
]):
    trow(sl, row, yx, Inches(3.16)+i*Inches(0.46), bg=WHITE if i%2==0 else GRAY_50)
txt(sl, "→ 정상: YOLO11n 정확도 우선\n   발열 시: YOLO26n으로 FPS 확보",
    Inches(0.85), Inches(5.4), Inches(5.5), Inches(0.75), size=13, color=TEXT_MID)

# 오른쪽
card(sl, Inches(7.05), Inches(2.05), Inches(5.9), Inches(5.2))
txt(sl, "Tracking: ByteTrack 선택", Inches(7.3), Inches(2.18),
    Inches(5.4), Inches(0.42), size=16, bold=True, color=INDIGO_900)
bx = [Inches(7.1), Inches(9.3), Inches(11.2)]
thdr(sl, ["항목", "SORT", "ByteTrack"], bx, Inches(2.72))
for i, row in enumerate([
    ("ID Switch",   "많음", "적음  ✓"),
    ("가려짐 처리", "불가", "가능  ✓"),
    ("속도",        "빠름", "빠름  ✓"),
]):
    trow(sl, row, bx, Inches(3.16)+i*Inches(0.46), bg=WHITE if i%2==0 else GRAY_50)
txt(sl, "→ 보행 중 가려짐(occlusion) 상황에서\n   객체 재추적이 반드시 필요",
    Inches(7.3), Inches(5.4), Inches(5.4), Inches(0.75), size=13, color=TEXT_MID)

callout(sl, '"YOLO11n 기반 탐지를 중심으로, 모바일 안정성을 위해 YOLO26n fallback 구조를 적용했습니다."')


# ════════════════════════════════════════════════════════════
# SLIDE 9 — 성능 지표
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 9, "핵심 AI 성능 지표", "수치로 말한다")

metrics = [
    ("XX FPS",  "모바일 추론 속도",  INDIGO,  INDIGO_50,  INDIGO_100),
    ("XX ms",   "평균 추론 시간",    GREEN,   GREEN_BG,   GREEN_200),
    ("XX%",     "mAP50 정확도",     PURPLE,  PURPLE_50,  PURPLE_200),
    ("XX ms",   "음성 응답 지연",    AMBER,   AMBER_BG,   AMBER_200),
]
for i, (val, label, ac, bg_c, bc) in enumerate(metrics):
    x = Inches(0.6) + i * Inches(3.18)
    rect(sl, x, Inches(2.05), Inches(2.95), Inches(1.72), fill=bg_c, lc=bc, lw=Pt(1))
    rect(sl, x, Inches(2.05), Inches(2.95), Inches(0.08), fill=ac, lc=None)
    txt(sl, val, x+Inches(0.1), Inches(2.2), Inches(2.75), Inches(0.82),
        size=38, bold=True, color=ac, align=PP_ALIGN.CENTER)
    txt(sl, label, x+Inches(0.1), Inches(3.05), Inches(2.75), Inches(0.35),
        size=12, color=TEXT_MID, align=PP_ALIGN.CENTER)

xs9 = [Inches(0.6), Inches(3.8), Inches(6.65), Inches(9.5), Inches(11.85)]
thdr(sl, ["항목", "YOLO11n (기본)", "YOLO26n (fallback)", "목표 기준"], xs9, Inches(4.1))
for i, row in enumerate([
    ("FPS (모바일)",     "[측정값] FPS",   "[측정값] FPS",   "10 FPS 이상"),
    ("mAP50",           "[측정값]%",      "[측정값]%",      "—"),
    ("평균 추론 시간",   "[측정값] ms",    "[측정값] ms",    "100ms 이하"),
    ("메모리 사용량",    "[측정값] MB",    "[측정값] MB",    "500MB 이하"),
    ("ID Switch 수",     "[측정값]회/분",  "—",              "낮을수록 좋음"),
]):
    trow(sl, row, xs9, Inches(4.54)+i*Inches(0.44), bg=WHITE if i%2==0 else GRAY_50, sz=12)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — 구현 및 실험 방법
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 10, "구현 및 실험 방법", "테스트 환경 및 비교 실험 설계")

card(sl, Inches(0.6), Inches(2.05), Inches(5.85), Inches(5.1))
rect(sl, Inches(0.6), Inches(2.05), Inches(5.85), Inches(0.52), fill=INDIGO, lc=None)
txt(sl, "테스트 환경", Inches(0.85), Inches(2.13), Inches(5.35), Inches(0.35),
    size=14, bold=True, color=WHITE)
for i, (k, v) in enumerate([
    ("스마트폰 기종",  "[기종명]"),
    ("Android 버전",   "[버전]"),
    ("CPU / GPU",      "[사양]"),
    ("테스트 장소",    "실내 복도, 외부 보도"),
    ("이동 조건",      "정지 상태 / 보행 중"),
    ("환경 조건",      "밝은 환경 / 저조도"),
]):
    y_ = Inches(2.72) + i * Inches(0.72)
    rect(sl, Inches(0.78), y_, Inches(5.5), Inches(0.62),
         fill=WHITE if i%2==0 else GRAY_50, lc=GRAY_200, lw=Pt(0.4))
    txt(sl, k, Inches(0.95), y_+Inches(0.12), Inches(2.2), Inches(0.35), size=13, color=TEXT_MID)
    txt(sl, v, Inches(3.25), y_+Inches(0.12), Inches(2.85), Inches(0.35), size=13, bold=True, color=TEXT)

card(sl, Inches(6.82), Inches(2.05), Inches(6.1), Inches(5.1))
rect(sl, Inches(6.82), Inches(2.05), Inches(6.1), Inches(0.52), fill=INDIGO, lc=None)
txt(sl, "비교 실험", Inches(7.07), Inches(2.13), Inches(5.6), Inches(0.35),
    size=14, bold=True, color=WHITE)
for i, (title, desc) in enumerate([
    ("모델 전환 구조",  "YOLO11n 단독  vs  YOLO11n + YOLO26n fallback"),
    ("Tracking 효과",  "ByteTrack 적용 전/후 ID Switch 수 비교"),
    ("환경 조건",       "밝은 환경 vs 저조도 탐지율 변화"),
    ("이동 조건",       "정지 상태 vs 보행 중 FPS 안정성 확인"),
]):
    y_ = Inches(2.72) + i * Inches(1.08)
    rect(sl, Inches(7.07), y_, Inches(0.07), Inches(0.52), fill=INDIGO_200, lc=None)
    txt(sl, title, Inches(7.3), y_, Inches(5.35), Inches(0.38),
        size=14, bold=True, color=INDIGO_900)
    txt(sl, desc, Inches(7.3), y_+Inches(0.4), Inches(5.35), Inches(0.55),
        size=12, color=TEXT_MID)


# ════════════════════════════════════════════════════════════
# SLIDE 11 — 성능 결과
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 11, "성능 결과", "실제로 얼마나 잘 되었는가")

xs11 = [Inches(0.6), Inches(3.7), Inches(6.05), Inches(8.3), Inches(10.6), Inches(12.73)]
thdr(sl, ["구성", "FPS", "mAP50", "안정성", "비고"], xs11, Inches(2.1))
for i, (row, bg_c) in enumerate([
    (("YOLO11n 단독",       "[측정값]", "높음", "발열 시 저하", "기본 모드"), WHITE),
    (("YOLO11n + fallback", "[측정값]", "중간", "안정적  ✓",    "권장 구성"), GREEN_BG),
]):
    trow(sl, row, xs11, Inches(2.54)+i*Inches(0.54), bg=bg_c)

for i, (title, desc) in enumerate([
    ("📊  FPS 시계열 그래프",    "fallback 전환 전/후 FPS 변화\n[ 그래프 이미지 삽입 ]"),
    ("📊  위험도 Smoothing 비교", "Smoothing 적용 전/후 비교\n[ 그래프 이미지 삽입 ]"),
]):
    x = Inches(0.6) + i * Inches(6.2)
    card(sl, x, Inches(3.65), Inches(5.9), Inches(3.5))
    rect(sl, x, Inches(3.65), Inches(5.9), Inches(0.5), fill=INDIGO, lc=None)
    txt(sl, title, x+Inches(0.2), Inches(3.72), Inches(5.5), Inches(0.35),
        size=13, bold=True, color=WHITE)
    txt(sl, desc, x+Inches(0.2), Inches(4.3), Inches(5.5), Inches(0.7),
        size=12, color=TEXT_MID)
    rect(sl, x+Inches(0.3), Inches(5.05), Inches(5.3), Inches(1.85),
         fill=GRAY_50, lc=GRAY_200, lw=Pt(0.5))
    txt(sl, "[ 그래프 이미지 삽입 ]",
        x+Inches(0.3), Inches(5.65), Inches(5.3), Inches(0.4),
        size=12, color=TEXT_LIGHT, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 12 — 진행 경과
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 12, "프로젝트 진행 경과", "주차별 타임라인")

timeline = [
    ("1주",  "아이디어\n기획",          GRAY_200,   TEXT_MID,    TEXT_MID),
    ("2주",  "모델 선정\nYOLO11n 확정", GRAY_200,   TEXT_MID,    TEXT_MID),
    ("3주",  "객체 탐지\n구현",          INDIGO_100, INDIGO_900,  INDIGO_900),
    ("4주",  "ByteTrack\n통합",          INDIGO_100, INDIGO_900,  INDIGO_900),
    ("5주",  "FPS 저하\n발견 ⚠",        RED_BG,     RED,         RED),
    ("6주",  "YOLO26n\nfallback 적용",   GREEN_BG,   GREEN,       GREEN),
    ("7주",  "서버 연동\n음성 피드백",    INDIGO_50,  INDIGO_900,  INDIGO_900),
    ("8주",  "통합 테스트\n데모 준비",    INDIGO,     WHITE,       WHITE),
]
TW = Inches(1.44)
BASE_X = Inches(0.5)
rect(sl, BASE_X+Inches(0.6), Inches(3.0), Inches(11.2), Pt(2), fill=GRAY_200, lc=None)

for i, (week, desc, node_bg, node_tc, desc_tc) in enumerate(timeline):
    x = BASE_X + i * Inches(1.55)
    cx = x + Inches(0.32)
    rect(sl, cx, Inches(2.78), Inches(0.46), Inches(0.46), fill=node_bg, lc=GRAY_200, lw=Pt(1))
    txt(sl, week, cx, Inches(2.8), Inches(0.46), Inches(0.34),
        size=11, bold=True, color=node_tc, align=PP_ALIGN.CENTER)
    lbl_y = Inches(2.05) if i%2==0 else Inches(3.38)
    rect(sl, x, lbl_y, TW, Inches(0.65), fill=node_bg, lc=GRAY_200, lw=Pt(0.6))
    txt(sl, desc, x+Inches(0.05), lbl_y+Inches(0.05), TW-Inches(0.08), Inches(0.58),
        size=11, color=desc_tc, align=PP_ALIGN.CENTER)

card(sl, Inches(0.6), Inches(4.6), Inches(12.13), Inches(1.2), bg=AMBER_BG, lc=AMBER_200)
rect(sl, Inches(0.6), Inches(4.6), Inches(0.09), Inches(1.2), fill=AMBER, lc=None)
txt(sl, "이슈 & 해결", Inches(0.87), Inches(4.7), Inches(2.5), Inches(0.38),
    size=14, bold=True, color=AMBER)
txt(sl, "5주차:  모바일 장시간 실행 시 발열로 YOLO11n FPS 저하 발생\n"
        "→  6주차: YOLO26n fallback 구조를 설계·적용하여 FPS 안정성 확보",
    Inches(0.87), Inches(5.1), Inches(11.65), Inches(0.6), size=13, color=TEXT)


# ════════════════════════════════════════════════════════════
# SLIDE 13 — 데모
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 13, "데모", "지금부터 직접 보여드립니다")

step_colors = [INDIGO, INDIGO, AMBER, GREEN, PURPLE]
steps = [
    ("01", "앱 실행",        "카메라 활성화, 초기 화면 확인"),
    ("02", "기본 탐지 확인", "YOLO11n 객체 탐지 동작"),
    ("03", "fallback 전환",  "부하 환경 → YOLO26n 자동 전환 확인"),
    ("04", "ID 추적 확인",   "ByteTrack Track ID 안정성"),
    ("05", "음성 경고 출력", "위험도 임계값 초과 → 음성 안내 출력"),
]
for i, ((num, title, desc), ac) in enumerate(zip(steps, step_colors)):
    y_ = Inches(2.12) + i * Inches(0.88)
    rect(sl, Inches(0.6), y_, Inches(0.52), Inches(0.58), fill=ac, lc=None)
    txt(sl, num, Inches(0.6), y_+Inches(0.1), Inches(0.52), Inches(0.38),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(1.28), y_+Inches(0.04), Inches(2.85), Inches(0.38),
        size=15, bold=True, color=TEXT)
    txt(sl, desc, Inches(1.28), y_+Inches(0.44), Inches(4.6), Inches(0.32),
        size=12, color=TEXT_MID)

card(sl, Inches(7.35), Inches(2.12), Inches(5.55), Inches(3.0))
rect(sl, Inches(7.35), Inches(2.12), Inches(5.55), Inches(0.5), fill=INDIGO, lc=None)
txt(sl, "시연 환경", Inches(7.6), Inches(2.19), Inches(5.05), Inches(0.35),
    size=14, bold=True, color=WHITE)
for i, (k, v) in enumerate([("스마트폰", "[기종명]"), ("장소", "실내 복도"), ("예상 시간", "5 ~ 7분")]):
    txt(sl, f"{k}  :  {v}", Inches(7.6), Inches(2.75)+i*Inches(0.58),
        Inches(5.05), Inches(0.42), size=14, color=TEXT)

card(sl, Inches(7.35), Inches(5.45), Inches(5.55), Inches(1.65), bg=GRAY_50, lc=GRAY_200)
txt(sl, "백업 영상 QR", Inches(7.6), Inches(5.6), Inches(5.05), Inches(0.38),
    size=13, bold=True, color=TEXT_MID)
txt(sl, "[ QR 코드 삽입 ]", Inches(7.6), Inches(6.1), Inches(5.05), Inches(0.45),
    size=13, color=TEXT_LIGHT, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 14 — 자체 평가 및 향후 방향
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
set_bg(sl)
hdr(sl, 14, "자체 평가 및 향후 방향", "객관적 평가와 다음 단계")

sections = [
    (GREEN,  GREEN_BG,  GREEN_200,  "잘 된 점", [
        "모바일 온디바이스 실시간 추론 성공",
        "YOLO26n fallback으로 FPS 안정성 확보",
        "ByteTrack으로 ID 안정성 향상",
        "사용자 중심 음성 피드백 설계",
    ]),
    (RED,    RED_BG,    RED_200,    "부족한 점", [
        "저조도 환경 탐지율 저하",
        "Depth 추정 정확도 한계",
        "실제 거리 캘리브레이션 미완성",
    ]),
    (INDIGO, INDIGO_50, INDIGO_100, "개선 방향", [
        "TensorRT 적용 → 추론 속도 향상",
        "경량 Depth 모델 교체 검토",
        "NPU 기기 대응 (Edge AI)",
        "실제 사용자 현장 테스트",
    ]),
]
for i, (ac, bg_c, bc, title, items) in enumerate(sections):
    x = Inches(0.55) + i * Inches(4.2)
    card(sl, x, Inches(2.05), Inches(3.95), Inches(5.1), bg=bg_c, lc=bc)
    rect(sl, x, Inches(2.05), Inches(3.95), Inches(0.1), fill=ac, lc=None)
    txt(sl, title, x+Inches(0.2), Inches(2.22), Inches(3.55), Inches(0.48),
        size=18, bold=True, color=ac)
    rect(sl, x+Inches(0.2), Inches(2.76), Inches(3.55), Pt(1.5), fill=ac, lc=None)
    for j, item in enumerate(items):
        txt(sl, f"·   {item}",
            x+Inches(0.2), Inches(2.92)+j*Inches(0.72), Inches(3.6), Inches(0.6),
            size=13, color=TEXT)


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
out = r"c:\VoiceGuide\VoiceGuide\docs\VoiceGuide_발표.pptx"
prs.save(out)
print("저장 완료:", out)
